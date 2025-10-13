#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SERP API Scraper for Educational Content

This script searches for educational content using SERP API and scrapes content from various sources.
It outputs scraped content to data/raw/serp_scrapes.jsonl for processing by the data pipeline.

Usage:
    python scripts/serp_scraper.py [subject]

Environment Variables:
    SERPAPI_KEY: Your SERP API key

Requirements:
    pip install google-search-results requests beautifulsoup4 pymupdf python-pptx python-docx python-dotenv
"""

import sys
import os
import json
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin
import time
import re
import uuid
from datetime import datetime, timezone

try:
    from serpapi import GoogleSearch
except ImportError:
    print(
        "Error: 'google-search-results' module not installed. Install with 'pip install google-search-results'."
    )
    sys.exit(1)

try:
    import fitz  # PyMuPDF for PDFs
except ImportError:
    print("Error: 'pymupdf' module not installed. Install with 'pip install pymupdf'.")
    sys.exit(1)

try:
    from pptx import Presentation  # For PPTX
except ImportError:
    print(
        "Error: 'python-pptx' module not installed. Install with 'pip install python-pptx'."
    )
    sys.exit(1)

try:
    from docx import Document  # For DOCX
except ImportError:
    print(
        "Error: 'python-docx' module not installed. Install with 'pip install python-docx'."
    )
    sys.exit(1)

try:
    from dotenv import load_dotenv

    load_dotenv()
except ImportError:
    print(
        "Warning: python-dotenv not installed. Install with 'pip install python-dotenv'."
    )

from zipfile import ZipFile  # Built-in for ZIP
import io  # For buffer handling
import tempfile

# HTTP headers for requests
headers = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
}


def is_video_url(url):
    """Check if URL points to video content."""
    video_extensions = [
        ".mp4",
        ".avi",
        ".mov",
        ".wmv",
        ".flv",
        ".webm",
        ".mkv",
        ".mpg",
        ".mpeg",
        ".3gp",
        ".ogv",
    ]
    video_domains = [
        "youtube.com",
        "vimeo.com",
        "dailymotion.com",
        "twitch.tv",
        "tiktok.com",
    ]
    url_lower = url.lower()
    if any(url_lower.endswith(ext) for ext in video_extensions):
        return True
    if any(domain in url_lower for domain in video_domains):
        return True
    return False


def is_academic_url(url):
    """Check if URL points to academic/course content that should be avoided."""
    # Whitelist: Allow these French sites that have good practical content
    whitelist_french = [
        "koor.fr",
        "python.doctor",
        "openclassrooms.com",
        "fun-mooc.fr",
        "commentcoder.com",
        "moncoachdata.com",
        "futureengineer.fr",
        "kinsta.com",  # Has French content
        "developpez.com",
        "linuxfr.org",
        "debian-facile.org",
        "ubuntu-fr.org",
    ]

    url_lower = url.lower()
    # Allow French sites not in the blocked list
    if any(site in url_lower for site in whitelist_french):
        return False

    academic_domains = [
        ".edu",
        ".ac.uk",
        ".ac.in",
        ".ac.au",
        ".ac.ca",
        "coursera.org",
        "udemy.com",
        "edx.org",
        "khanacademy.org",
        "university",
        "college",
        "geeksforgeeks.org",
        "w3schools.com",
        "tutorialspoint.com",
        "javatpoint.com",
        "programiz.com",
        "codecademy.com",
        "freecodecamp.org",
        "hackerrank.com",
        "leetcode.com",
        "stackoverflow.com",
        "github.com",
        "medium.com",
        "towardsdatascience.com",
        "realpython.com",
        "python.org",
        "docs.python.org",
        # Add more English-focused sites
        "datacamp.com",
        "kaggle.com",
        "pytorch.org",
        "tensorflow.org",
    ]
    result = any(domain in url_lower for domain in academic_domains)
    if result:
        print(f"DEBUG: Academic/known English site URL detected: {url}")
    return result


def scrape_content(url, is_sub_page=False):
    """Scrape content from various file types and web pages."""
    try:
        response = requests.get(url, headers=headers, timeout=(5, 30))
        response.raise_for_status()

        # Detect based on URL extension or Content-Type
        ext = url.lower().split(".")[-1] if "." in url else ""
        content_type = response.headers.get("Content-Type", "").lower()

        # Handle PPTX files
        if (
            ext == "pptx"
            or "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            in content_type
        ):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_file:
                temp_file.write(response.content)
                temp_path = temp_file.name
            prs = Presentation(temp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            os.remove(temp_path)
            title = os.path.basename(url)
            return text, title, None  # Full text

        # Handle PDF files
        elif ext == "pdf" or "application/pdf" in content_type:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                temp_file.write(response.content)
                temp_path = temp_file.name
            doc = fitz.open(temp_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            os.remove(temp_path)
            title = os.path.basename(url)
            return text, title, None  # Full text

        # Handle DOCX files
        elif (
            ext == "docx"
            or "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            in content_type
        ):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_file:
                temp_file.write(response.content)
                temp_path = temp_file.name
            doc = Document(temp_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            os.remove(temp_path)
            title = os.path.basename(url)
            return text, title, None  # Full text

        # Handle ZIP files
        elif ext == "zip" or "application/zip" in content_type:
            with io.BytesIO(response.content) as zip_buffer:
                with ZipFile(zip_buffer) as z:
                    text = ""
                    for file_name in z.namelist():
                        if file_name.lower().endswith(
                            (".txt", ".md", ".html")
                        ):  # Extract text from simple files in ZIP
                            with z.open(file_name) as f:
                                text += (
                                    f.read().decode("utf-8", errors="ignore") + "\n\n"
                                )
            title = os.path.basename(url)
            return text, title, None  # Full text

        # Handle TXT or plain text
        elif ext == "txt" or "text/plain" in content_type:
            text = response.text
            title = os.path.basename(url) or "Untitled"
            return text, title, None  # Full text

        # Default: HTML
        else:
            response.encoding = "utf-8"
            soup = BeautifulSoup(response.text, "html.parser")
            text = soup.get_text(separator="\n", strip=True)
            title = soup.title.string if soup.title else "Untitled"
            return text, title, soup  # Full text

    except Exception as e:
        return f"Error scraping {url}: {str(e)}", "Error", None


def extract_relevant_links(soup, base_url):
    """Extract relevant educational links from HTML content."""
    if not soup:
        return []
    links = []
    keywords = [
        "course",
        "td",
        "tp",
        "exercise",
        "lesson",
        "tutorial",
        "assignment",
        "homework",
    ]
    for a in soup.find_all("a", href=True):
        href = a["href"]
        if not (href.startswith("http://") or href.startswith("https://")):
            href = urljoin(base_url, href)
            if not (href.startswith("http://") or href.startswith("https://")):
                continue
        link_text = a.text.strip()
        if any(keyword in link_text.lower() for keyword in keywords):
            links.append((link_text, href))
    return links[:10]  # Limit to 10 relevant links


def clean_filename(text):
    """Clean text for use as filename."""
    # Remove or replace invalid characters for Windows filenames
    text = text.lower()
    # Replace accented characters
    text = (
        text.replace("é", "e")
        .replace("è", "e")
        .replace("ê", "e")
        .replace("à", "a")
        .replace("ç", "c")
    )
    # Remove forbidden characters and replace with '_'
    text = re.sub(r'[<>:"/\\|?*\n\r\t]', "_", text)
    # Replace spaces and special characters with '_'
    text = re.sub(r"\s+", "_", text)
    # Replace multiple '_' with single '_'
    text = re.sub(r"_+", "_", text)
    # Remove '_' from start or end
    text = text.strip("_")
    # Limit length to 100 characters
    text = text[:100]
    # If empty after cleaning, use 'unknown'
    return text if text else "unknown"


def save_to_jsonl(subject, url, content, title, source_type="main", link_name=None):
    """Save scraped content to JSONL file for pipeline processing."""
    try:
        # Ensure output directory exists
        output_dir = "data/raw"
        os.makedirs(output_dir, exist_ok=True)
        output_file = os.path.join(output_dir, "serp_scrapes.jsonl")

        # Create record
        record = {
            "id": str(uuid.uuid4()),
            "subject": subject,
            "url": url,
            "title": title,
            "text": content,
            "source_type": source_type,  # "main" or "sub_page"
            "link_name": link_name,
            "scraped_at": datetime.now(timezone.utc).isoformat(),
            "word_count": len(content.split()) if content else 0,
        }

        # Append to JSONL file
        with open(output_file, "a", encoding="utf-8") as f:
            json.dump(record, f, ensure_ascii=False)
            f.write("\n")

        return output_file, None
    except Exception as e:
        return None, f"Error saving {url}: {str(e)}"


def main():
    """Main scraping function."""
    errors = []  # List to store errors

    # Get subject from command line or user input
    if len(sys.argv) > 1:
        subject = sys.argv[1]
    else:
        subject = input("Enter the subject name: ")

    # Get API key from environment
    api_key = os.getenv("SERPAPI_KEY")
    if not api_key:
        print(
            "Error: SERPAPI_KEY not found. Set it in a .env file or environment variable."
        )
        sys.exit(1)

    # Create search query for actual educational content (tutorials, guides, docs)
    # Focus on French websites with practical programming examples, TP, TD, exercises
    # Avoid known English sites, focus on lesser-known French sources
    content_queries = [
        f"tutoriel {subject} programmation exemples code pratique",
        f"apprendre {subject} exemples pratiques guide débutant",
        f"{subject} tutoriel codage avec exemples fonctionnels",
        f"exercices pratiques {subject} programmation solutions",
        f"tutoriel {subject} étape par étape code détaillé",
        f"exemples {subject} programmation expliqués en détail",
        f"cours {subject} programmation TP TD exercices",
        f"{subject} développement tutoriel code source",
        f"{subject} programmation tutoriel français exemples",
        f"guide {subject} apprentissage pratique code",
        f"{subject} exercices corrigés programmation",
        f"tutoriel {subject} développement avec code",
        # New queries specifically for TD/TP and exercises
        f"TD {subject} travaux dirigés exercices corrigés",
        f"TP {subject} travaux pratiques programmation",
        f"exercices {subject} corrigés solutions détaillées",
        f"problèmes {subject} résolution avec code",
        f"ateliers {subject} programmation pratique",
        f"laboratoire {subject} exercices informatiques",
        f"pratique {subject} coding exercises français",
        f"{subject} algorithmique exercices TD TP",
        f"programmation {subject} exercices niveau débutant intermédiaire",
        f"corrigés {subject} exercices programmation Python",
        # Site-specific searches for known good French sites
        f"site:koor.fr {subject} exercices TP TD",
        f"site:python.doctor {subject} tutoriel exercices",
        f"site:openclassrooms.com {subject} TP exercices",
        f"site:fun-mooc.fr {subject} travaux pratiques",
        f"site:commentcoder.com {subject} exercices corrigés",
    ]

    # Use multiple queries to get diverse content sources
    all_results = []
    for query in content_queries:
        print(f"Searching for: {query}")
        try:
            start = 0
            while (
                len(all_results) < 50
            ):  # Get more results per query for better coverage
                params = {
                    "q": query,
                    "hl": "fr",
                    "gl": "fr",
                    "num": 20,  # Increase to 20 results per page
                    "start": start,
                    "api_key": api_key,
                }
                search = GoogleSearch(params)
                results = search.get_dict()

                if "error" in results:
                    print(f"SERP API Error: {results['error']}")
                    break

                organic_results = results.get("organic_results", [])
                print(
                    f"Found {len(organic_results)} organic results for query: {query}"
                )
                if not organic_results:
                    break

                all_results.extend(organic_results)
                start += 10

        except Exception as e:
            print(f"Error with query '{query}': {e}")
            continue

    print(f"Total results collected: {len(all_results)}")

    if not all_results:
        print("No results found. Try a different subject or check your connection.")
        return

    # Process each result
    for i, result in enumerate(all_results[:20], start=1):
        url = result.get("link")
        print(f"Processing URL {i}: {url}")
        if is_video_url(url):
            print(f"Skipping video URL: {url}")
            continue
        if is_academic_url(url):
            print(f"Skipping academic/course URL: {url}")
            continue

        print(f"\n--- Site {i}: {url} ---")
        try:
            content, title, soup = scrape_content(url)
            if "Error scraping" in content:
                errors.append(f"Error in Site {i}: {content}")
                print(f"Error in Site {i}: {content}")
                continue

            # Save main content
            filepath, save_error = save_to_jsonl(subject, url, content, title, "main")
            if save_error:
                errors.append(f"Error in Site {i}: {save_error}")
                print(f"Error in Site {i}: {save_error}")
                continue

            print(f"Main content saved to: {filepath}")
            print(content[:200] + "..." if len(content) > 200 else content)

            # Extract and scrape relevant sub-pages if HTML
            links = extract_relevant_links(soup, url)
            if links:
                print(f"Found {len(links)} relevant links. Scraping sub-pages...")
                for j, (link_text, sub_url) in enumerate(links, start=1):
                    if is_video_url(sub_url):
                        print(f"Skipping video URL: {sub_url}")
                        continue

                    try:
                        time.sleep(1)  # Delay to avoid server overload
                        sub_content, sub_title, _ = scrape_content(
                            sub_url, is_sub_page=True
                        )
                        if "Error scraping" in sub_content:
                            errors.append(
                                f"Error in Link {j} of Site {i} ({link_text}): {sub_content}"
                            )
                            print(f"Error in Link {j} of Site {i}: {sub_content}")
                            continue

                        link_name = (
                            clean_filename(link_text)
                            or os.path.basename(sub_url).split(".")[0]
                        )
                        sub_filepath, sub_save_error = save_to_jsonl(
                            subject,
                            sub_url,
                            sub_content,
                            sub_title,
                            "sub_page",
                            link_name,
                        )
                        if sub_save_error:
                            errors.append(
                                f"Error in Link {j} of Site {i} ({link_text}): {sub_save_error}"
                            )
                            print(f"Error in Link {j} of Site {i}: {sub_save_error}")
                            continue

                        print(f"  - Sub-page {j} ({link_text}): {sub_url}")
                        print(f"    Content saved to: {sub_filepath}")

                        # Extract and scrape sub-sub-pages from this sub-page
                        try:
                            sub_soup = BeautifulSoup(sub_content, "html.parser")
                            sub_links = extract_relevant_links(sub_soup, sub_url)
                            if sub_links:
                                print(
                                    f"    Found {len(sub_links)} sub-links. Scraping sub-sub-pages..."
                                )
                                for k, (sub_link_text, sub_sub_url) in enumerate(
                                    sub_links[:5], start=1
                                ):  # Limit to 5 sub-sub-links
                                    if is_video_url(sub_sub_url):
                                        print(f"    Skipping video URL: {sub_sub_url}")
                                        continue
                                    if is_academic_url(sub_sub_url):
                                        print(
                                            f"    Skipping academic URL: {sub_sub_url}"
                                        )
                                        continue

                                    try:
                                        time.sleep(1)  # Delay to avoid server overload
                                        sub_sub_content, sub_sub_title, _ = (
                                            scrape_content(
                                                sub_sub_url, is_sub_page=True
                                            )
                                        )
                                        if "Error scraping" in sub_sub_content:
                                            errors.append(
                                                f"Error in Sub-link {k} of Sub-page {j} ({sub_link_text}): {sub_sub_content}"
                                            )
                                            print(
                                                f"    Error in Sub-link {k} of Sub-page {j}: {sub_sub_content}"
                                            )
                                            continue

                                        sub_link_name = (
                                            clean_filename(sub_link_text)
                                            or os.path.basename(sub_sub_url).split(".")[
                                                0
                                            ]
                                        )
                                        sub_sub_filepath, sub_sub_save_error = (
                                            save_to_jsonl(
                                                subject,
                                                sub_sub_url,
                                                sub_sub_content,
                                                sub_sub_title,
                                                "sub_sub_page",
                                                sub_link_name,
                                            )
                                        )
                                        if sub_sub_save_error:
                                            errors.append(
                                                f"Error in Sub-link {k} of Sub-page {j} ({sub_link_text}): {sub_sub_save_error}"
                                            )
                                            print(
                                                f"    Error in Sub-link {k} of Sub-page {j}: {sub_sub_save_error}"
                                            )
                                            continue

                                        print(
                                            f"      - Sub-sub-page {k} ({sub_link_text}): {sub_sub_url}"
                                        )
                                        print(
                                            f"        Content saved to: {sub_sub_filepath}"
                                        )

                                    except Exception as sub_sub_e:
                                        error_msg = f"Error in Sub-link {k} of Sub-page {j}: {str(sub_sub_e)}"
                                        errors.append(error_msg)
                                        print(f"    {error_msg}")
                                        continue

                        except Exception as sub_links_e:
                            print(
                                f"    Error extracting sub-links from sub-page {j}: {str(sub_links_e)}"
                            )
                            continue
                        print(
                            sub_content[:200] + "..."
                            if len(sub_content) > 200
                            else sub_content
                        )

                    except Exception as e:
                        error_msg = (
                            f"Error in Link {j} of Site {i} ({link_text}): {str(e)}"
                        )
                        errors.append(error_msg)
                        print(error_msg)
                        continue

                print("\n" + "=" * 80 + "\n")

        except Exception as e:
            error_msg = f"Error in Site {i}: {str(e)}"
            errors.append(error_msg)
            print(error_msg)
            continue

    # Display all collected errors
    if errors:
        print("\n=== Error Summary ===")
        for error in errors:
            print(error)
    else:
        print("\nNo errors encountered.")


if __name__ == "__main__":
    main()
